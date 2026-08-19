#!/usr/bin/env python3
"""Build the editable Reviewer-routing and recovery mechanism figure."""

from __future__ import annotations

import hashlib
import json
import shutil
import subprocess
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

REPORT = Path(__file__).resolve().parents[1]
FIGURES = REPORT / "figures"
EVIDENCE = REPORT / "evidence" / "swebench_pro"
STATS_PATH = EVIDENCE / "reviewer_mechanism_stats.json"
TASKS_PATH = EVIDENCE / "reviewer_interventions.csv"
PPTX_PATH = FIGURES / "reviewer_mechanism.pptx"
PDF_PATH = FIGURES / "reviewer_mechanism.pdf"
SVG_PATH = FIGURES / "reviewer_mechanism.svg"
PNG_PATH = FIGURES / "reviewer_mechanism.png"
MACROS_PATH = FIGURES / "reviewer_metrics.tex"
PROVENANCE_PATH = FIGURES / "reviewer_mechanism.provenance.json"
ANIME = FIGURES / "assets" / "anime"
ENGINEER_AVATAR = ANIME / "engineer.png"
REVIEWER_AVATAR = ANIME / "reviewer.png"
MOUNTAIN_STRIP = ANIME / "mountain_strip.png"


WHITE = RGBColor(255, 253, 248)
PAPER = RGBColor(251, 247, 238)
INK = RGBColor(36, 70, 93)
MUTED = RGBColor(102, 113, 125)
BLUE = RGBColor(49, 91, 206)
DEEP = RGBColor(23, 59, 112)
PALE_BLUE = RGBColor(236, 241, 253)
GOLD = RGBColor(195, 138, 32)
PALE_GOLD = RGBColor(252, 246, 229)
GRAY = RGBColor(122, 131, 142)
LIGHT_GRAY = RGBColor(217, 224, 231)
PANEL = RGBColor(255, 253, 248)


def sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 10,
    color: RGBColor = INK,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> Any:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = frame.margin_right = Inches(0.02)
    frame.margin_top = frame.margin_bottom = Inches(0.01)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_rect(slide, x: float, y: float, w: float, h: float, fill: RGBColor, line: RGBColor = LIGHT_GRAY) -> Any:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(0.8)
    shape.shadow.inherit = False
    return shape


def add_line(slide, x1: float, y1: float, x2: float, y2: float, color: RGBColor, width: float = 1.5) -> Any:
    shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    shape.line.color.rgb = color
    shape.line.width = Pt(width)
    return shape


def add_metric_box(slide, x: float, y: float, w: float, h: float, value: str, label: str, *, fill: RGBColor, color: RGBColor = DEEP) -> None:
    add_rect(slide, x, y, w, h, fill)
    add_text(slide, value, x + 0.10, y + 0.08, w - 0.20, h * 0.47, size=18, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(slide, label, x + 0.10, y + h * 0.52, w - 0.20, h * 0.35, size=9.3, color=MUTED, align=PP_ALIGN.CENTER)


def write_macros(stats: dict[str, Any]) -> None:
    routing = stats["routing"]
    outcomes = stats["external_reviewer_outcomes"]
    workload = stats["descriptive_workload"]
    values = {
        "ReviewerTasks": stats["tasks"],
        "ReviewerInvoked": routing["external_reviewer_tasks"],
        "ReviewerInvokedPercent": f"{100 * routing['external_reviewer_share']:.1f}",
        "ReviewerSkipped": routing["engineer_self_review_tasks"],
        "ReviewerSkippedPercent": f"{100 * routing['engineer_self_review_share']:.1f}",
        "ReviewerFirstAccepted": outcomes["first_review_accepted"],
        "ReviewerFirstBlocked": outcomes["first_review_blocked"],
        "ReviewerRevisionRequested": outcomes["revision_requested"],
        "ReviewerVerifierRecovered": outcomes["official_verifier_resolved_after_revision"],
        "ReviewerStrictRescues": outcomes["reviewer_accepted_after_revision"],
        "ReviewerTokenRatio": f"{workload['external_reviewer_mean_solve_input_tokens'] / workload['self_review_mean_solve_input_tokens']:.2f}",
        "ReviewerTimeRatio": f"{workload['external_reviewer_mean_active_seconds'] / workload['self_review_mean_active_seconds']:.2f}",
    }
    MACROS_PATH.write_text(
        "".join(f"\\newcommand{{\\{key}}}{{{value}}}\n" for key, value in values.items()),
        encoding="utf-8",
    )


def main() -> int:
    stats = json.loads(STATS_PATH.read_text(encoding="utf-8"))
    write_macros(stats)
    routing = stats["routing"]
    outcomes = stats["external_reviewer_outcomes"]
    workload = stats["descriptive_workload"]

    prs = Presentation()
    prs.slide_width = Inches(12)
    prs.slide_height = Inches(2.65)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    slide.shapes.add_picture(str(MOUNTAIN_STRIP), Inches(0), Inches(2.44), width=Inches(12), height=Inches(0.21))
    add_rect(slide, 0, 0, 12, 0.045, INK, INK)
    # Routing split.
    x, y, w, h = 0.35, 0.16, 3.35, 2.15
    add_rect(slide, x, y, w, h, PANEL)
    slide.shapes.add_picture(str(ENGINEER_AVATAR), Inches(x + 2.60), Inches(y + 0.02), width=Inches(0.58), height=Inches(0.58))
    add_text(slide, "(a) Routing", x + 0.16, y + 0.08, 2.3, 0.28, size=11.5, bold=True, color=DEEP)
    bar_x, bar_y, bar_w = x + 0.18, y + 0.58, w - 0.36
    reviewer_share = float(routing["external_reviewer_share"])
    add_rect(slide, bar_x, bar_y, bar_w, 0.30, WHITE)
    add_rect(slide, bar_x, bar_y, bar_w * reviewer_share, 0.30, BLUE, BLUE)
    add_rect(slide, bar_x + bar_w * reviewer_share, bar_y, bar_w * (1 - reviewer_share), 0.30, GRAY, GRAY)
    add_text(slide, f"{routing['external_reviewer_tasks']} · {100 * reviewer_share:.1f}%", bar_x + 0.06, bar_y, bar_w * reviewer_share - 0.10, 0.30, size=9.2, color=WHITE, bold=True)
    add_text(slide, f"{routing['engineer_self_review_tasks']} · {100 * (1-reviewer_share):.1f}%", bar_x + bar_w * reviewer_share, bar_y, bar_w * (1-reviewer_share) - 0.04, 0.30, size=8.6, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Reviewer", bar_x, bar_y + 0.38, 1.45, 0.24, size=9.5, bold=True, color=BLUE)
    add_text(slide, "Self-review", bar_x + 1.62, bar_y + 0.38, 1.30, 0.24, size=9.5, bold=True, color=GRAY, align=PP_ALIGN.RIGHT)
    token_ratio = workload["external_reviewer_mean_solve_input_tokens"] / workload["self_review_mean_solve_input_tokens"]
    time_ratio = workload["external_reviewer_mean_active_seconds"] / workload["self_review_mean_active_seconds"]
    add_text(slide, f"{token_ratio:.2f}× tokens · {time_ratio:.2f}× time", x + 0.18, y + 1.56, w - 0.36, 0.28, size=11, color=DEEP, bold=True, align=PP_ALIGN.CENTER)

    # Intervention funnel.
    x, y, w, h = 3.92, 0.16, 7.73, 2.15
    add_rect(slide, x, y, w, h, WHITE)
    slide.shapes.add_picture(str(REVIEWER_AVATAR), Inches(x + 6.93), Inches(y + 0.02), width=Inches(0.58), height=Inches(0.58))
    add_text(slide, "(b) Revision recovery", x + 0.16, y + 0.08, 5.2, 0.28, size=11.5, bold=True, color=DEEP)
    add_metric_box(slide, x + 0.18, y + 0.60, 1.30, 1.02, str(routing["external_reviewer_tasks"]), "Invoked", fill=PALE_BLUE)
    add_line(slide, x + 1.52, y + 1.10, x + 2.02, y + 0.76, LIGHT_GRAY, 1.5)
    add_line(slide, x + 1.52, y + 1.10, x + 2.02, y + 1.56, GOLD, 1.8)
    add_metric_box(slide, x + 2.06, y + 0.40, 1.32, 0.66, str(outcomes["first_review_accepted"]), "Accepted", fill=PALE_BLUE)
    add_metric_box(slide, x + 2.06, y + 1.24, 1.32, 0.66, str(outcomes["revision_requested"]), "Revise", fill=PALE_GOLD, color=GOLD)
    add_line(slide, x + 3.42, y + 1.57, x + 3.92, y + 1.57, GOLD, 1.8)
    add_metric_box(slide, x + 3.96, y + 1.24, 1.42, 0.66, str(outcomes["official_verifier_resolved_after_revision"]), "Verifier pass", fill=PALE_BLUE)
    add_line(slide, x + 5.42, y + 1.57, x + 5.88, y + 1.57, GOLD, 1.8)
    add_metric_box(slide, x + 5.92, y + 1.24, 1.46, 0.66, str(outcomes["reviewer_accepted_after_revision"]), "Strict rescue", fill=PALE_GOLD, color=GOLD)
    prs.save(PPTX_PATH)

    libreoffice = shutil.which("libreoffice") or shutil.which("soffice")
    if not libreoffice:
        raise SystemExit("LibreOffice is required for PDF export")
    if PDF_PATH.exists():
        PDF_PATH.unlink()
    completed = subprocess.run([libreoffice, "--headless", "--convert-to", "pdf", "--outdir", str(FIGURES), str(PPTX_PATH)], capture_output=True, text=True)
    if completed.returncode != 0 or not PDF_PATH.is_file():
        raise SystemExit(f"Reviewer figure export failed: {completed.stderr}")
    ghostscript = shutil.which("gs")
    if ghostscript:
        compatible = FIGURES / "reviewer_mechanism.compat.pdf"
        subprocess.run([ghostscript, "-q", "-dNOPAUSE", "-dBATCH", "-sDEVICE=pdfwrite", "-dCompatibilityLevel=1.5", f"-sOutputFile={compatible}", str(PDF_PATH)], check=True)
        compatible.replace(PDF_PATH)
    subprocess.run(["pdftocairo", "-svg", str(PDF_PATH), str(SVG_PATH)], check=True)
    subprocess.run(["pdftoppm", "-singlefile", "-png", "-r", "180", str(PDF_PATH), str(PNG_PATH.with_suffix(""))], check=True)

    provenance = {
        "figure_id": "reviewer-mechanism",
        "reader_question": "How often is an independent Reviewer invoked, and how many revision-requested tasks are recovered?",
        "claim": "Reviewer is invoked on 466 of 731 tasks; 43 receive revision requests, 34 later pass the official verifier, and 22 complete the strict review-loop rescue.",
        "scope": "Adaptive routing analysis, not randomized ablation.",
        "visual_style": "anime research field note with cream paper, navy linework, mountain-ridge motif, and shared Engineer/Reviewer characters; flow counts remain exact",
        "character_assets": [str(ENGINEER_AVATAR.relative_to(FIGURES)), str(REVIEWER_AVATAR.relative_to(FIGURES))],
        "inputs": {STATS_PATH.name: sha256(STATS_PATH), TASKS_PATH.name: sha256(TASKS_PATH)},
        "outputs": {path.name: sha256(path) for path in (PPTX_PATH, PDF_PATH, SVG_PATH, PNG_PATH, MACROS_PATH)},
        "editable_source": PPTX_PATH.name,
    }
    PROVENANCE_PATH.write_text(json.dumps(provenance, indent=2, sort_keys=True) + "\n", encoding="utf-8")
    print(PPTX_PATH)
    print(PDF_PATH)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
